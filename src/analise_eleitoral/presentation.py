from __future__ import annotations
from pathlib import Path
from typing import Any
from pptx import Presentation
from pptx.util import Inches, Pt


def _title(slide,text:str)->None:
    box=slide.shapes.add_textbox(Inches(.7),Inches(.45),Inches(12),Inches(.7)); p=box.text_frame.paragraphs[0]; p.text=text; p.font.size=Pt(30); p.font.bold=True

def _bullets(slide,items:list[str])->None:
    box=slide.shapes.add_textbox(Inches(.9),Inches(1.5),Inches(11.4),Inches(5.2)); tf=box.text_frame
    for i,text in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=text; p.font.size=Pt(19); p.space_after=Pt(12)

def build_presentation(snapshot:dict[str,Any],chart_dir:Path,out_path:Path)->None:
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5); blank=prs.slide_layouts[6]
    s=prs.slides.add_slide(blank); _title(s,'Análise Eleitoral 2026'); _bullets(s,[f"Snapshot {snapshot['meta']['snapshot_date']} · v{snapshot['meta']['pipeline_version']}",'Análise descritiva, apartidária e reprodutível.'])
    s=prs.slides.add_slide(blank); _title(s,'Visão geral'); _bullets(s,[f"{snapshot['candidaturas']['candidaturas_total']:,} registros de candidatura",f"Receitas agregadas: R$ {snapshot['financas']['total_receitas']:,.2f}",f"Despesas contratadas: R$ {snapshot['financas']['total_despesas_contratadas']:,.2f}"])
    for title,name in [('Candidaturas por cargo','candidaturas_por_cargo.png'),('Receitas por fonte','receitas_por_fonte.png'),('Despesas contratadas','despesas_top10.png')]:
        s=prs.slides.add_slide(blank); _title(s,title); p=chart_dir/name
        if p.exists(): s.shapes.add_picture(str(p),Inches(.8),Inches(1.35),width=Inches(11.8))
    s=prs.slides.add_slide(blank); _title(s,'Qualidade e governança'); gate=snapshot.get('qualidade',{}).get('publication_gate',{}); _bullets(s,[f"Gate de publicação: {'APROVADO' if gate.get('ok') else 'não registrado/reprovado'}",'Freshness e integridade bloqueiam cargas inválidas.','Snapshots reprovados são enviados à quarentena e não substituem o latest válido.'])
    out_path.parent.mkdir(parents=True,exist_ok=True); prs.save(out_path)
